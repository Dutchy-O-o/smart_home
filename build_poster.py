"""
Builds an A1 portrait (594 x 841 mm) project poster following the OrnekPoster.pptx
guidelines. Output: SmartHome_Poster.pptx in the project root.
"""

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os


# ---------------------------------------------------------------------------
# Theme palette (DEU navy + accent)
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x0B, 0x2A, 0x5B)   # main brand
NAVY_DK = RGBColor(0x07, 0x1B, 0x3D)
ACCENT = RGBColor(0xE3, 0x6F, 0x1E)  # warm orange (mood/AI accent)
ACCENT2 = RGBColor(0x1D, 0xB9, 0x54)  # spotify green
TEAL = RGBColor(0x0F, 0x8E, 0xA8)
INK = RGBColor(0x12, 0x18, 0x2A)
SUB = RGBColor(0x44, 0x4C, 0x66)
LINE = RGBColor(0xC9, 0xCF, 0xDC)
BG_SOFT = RGBColor(0xF4, 0xF6, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHIP_BG = RGBColor(0xEA, 0xEF, 0xFA)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill, line=None, line_w=None,
             corner=False, shadow=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE,
        x, y, w, h
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w is not None:
            shp.line.width = line_w
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_paragraphs(slide, x, y, w, h, paragraphs, size=14, color=INK,
                   align=PP_ALIGN.LEFT, font="Calibri", line_spacing=1.1):
    """paragraphs: list of dicts {text, bold?, size?, color?, bullet?}"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.15)
    tf.margin_right = Cm(0.15)
    tf.margin_top = Cm(0.1)
    tf.margin_bottom = Cm(0.1)
    for i, item in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = item.get("align", align)
        p.line_spacing = item.get("line_spacing", line_spacing)
        if item.get("bullet"):
            # add bullet via XML
            pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
            buChar = etree.SubElement(pPr, qn("a:buChar"))
            buChar.set("char", "•")
            pPr.set("indent", "-228600")
            pPr.set("marL", "228600")
        runs = item.get("runs")
        if runs is None:
            runs = [{"text": item["text"],
                     "bold": item.get("bold", False),
                     "size": item.get("size", size),
                     "color": item.get("color", color),
                     "italic": item.get("italic", False)}]
        for k, run_spec in enumerate(runs):
            r = p.add_run()
            r.text = run_spec["text"]
            r.font.name = font
            r.font.size = Pt(run_spec.get("size", size))
            r.font.bold = run_spec.get("bold", False)
            r.font.italic = run_spec.get("italic", False)
            r.font.color.rgb = run_spec.get("color", color)
    return tb


def add_section(slide, x, y, w, h, title, body_paragraphs,
                title_color=WHITE, accent=NAVY, body_size=12.5,
                title_size=20):
    """Card with colored title bar + white body."""
    # card background
    card = add_rect(slide, x, y, w, h, WHITE, line=LINE, line_w=Pt(0.75),
                    corner=True)
    # title bar
    title_h = Cm(1.6)
    bar = add_rect(slide, x, y, w, title_h, accent, corner=True)
    add_text(slide, x + Cm(0.4), y, w - Cm(0.6), title_h, title,
             size=title_size, bold=True, color=title_color,
             anchor=MSO_ANCHOR.MIDDLE)
    # body
    add_paragraphs(slide, x + Cm(0.2), y + title_h + Cm(0.1),
                   w - Cm(0.4), h - title_h - Cm(0.2),
                   body_paragraphs, size=body_size)


def add_chip(slide, x, y, label, fill=CHIP_BG, color=NAVY, size=11):
    w = Cm(0.45 * len(label) + 0.8)
    h = Cm(0.95)
    chip = add_rect(slide, x, y, w, h, fill, corner=True)
    add_text(slide, x, y, w, h, label, size=size, bold=True,
             color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return w


# ---------------------------------------------------------------------------
# Build the presentation
# ---------------------------------------------------------------------------
prs = Presentation()
# A1 portrait 594 x 841 mm
prs.slide_width = Cm(59.4)
prs.slide_height = Cm(84.1)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# -------- Page background --------
add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, BG_SOFT)

# ============================================================
# HEADER (0 .. 11 cm)
# ============================================================
HEAD_H = Cm(11.0)
add_rect(slide, 0, 0, prs.slide_width, HEAD_H, NAVY)
add_rect(slide, 0, HEAD_H, prs.slide_width, Cm(0.25), ACCENT)

# University block (top-left)
add_text(slide, Cm(1.0), Cm(0.6), Cm(20), Cm(1.4),
         "DOKUZ EYLÜL ÜNİVERSİTESİ",
         size=22, bold=True, color=WHITE)
add_text(slide, Cm(1.0), Cm(2.0), Cm(28), Cm(1.0),
         "Mühendislik Fakültesi  ·  Bilgisayar Mühendisliği Bölümü",
         size=15, color=RGBColor(0xCF, 0xDA, 0xF0))
add_text(slide, Cm(1.0), Cm(2.9), Cm(28), Cm(0.9),
         "Bitirme Projesi  ·  Ekim 2025  ·  İzmir",
         size=13, italic=True, color=RGBColor(0xCF, 0xDA, 0xF0))

# Project number (top-right)
pno_w = Cm(13.5)
pno_x = prs.slide_width - pno_w - Cm(1.0)
add_rect(slide, pno_x, Cm(0.7), pno_w, Cm(3.0), WHITE, corner=True)
add_text(slide, pno_x, Cm(0.85), pno_w, Cm(0.9),
         "Proje No",
         size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(slide, pno_x, Cm(1.7), pno_w, Cm(1.0),
         "2022510085 · 105 · 111",
         size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
add_text(slide, pno_x, Cm(2.7), pno_w, Cm(0.8),
         "Bilgisayar Mühendisliği · 2025–2026",
         size=11, color=SUB, align=PP_ALIGN.CENTER)

# Title
add_text(slide, Cm(1.0), Cm(4.4), prs.slide_width - Cm(2.0), Cm(2.3),
         "IoT Tabanlı Akıllı Ev Sistemi:",
         size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Cm(1.0), Cm(6.6), prs.slide_width - Cm(2.0), Cm(2.0),
         "Duygu-Odaklı Kişiselleştirme ile Adaptif Yaşam Alanı",
         size=34, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
         italic=True)

# Team strip
team_y = Cm(8.7)
add_text(slide, Cm(1.0), team_y, prs.slide_width - Cm(2.0), Cm(1.0),
         "Emre Akkaya  ·  Ramazan Denli  ·  Boran Bereketli",
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Cm(1.0), team_y + Cm(1.0), prs.slide_width - Cm(2.0),
         Cm(1.0),
         "Danışman:  Dr. Öğr. Üyesi Yunus DOĞAN",
         size=16, color=RGBColor(0xCF, 0xDA, 0xF0),
         align=PP_ALIGN.CENTER, italic=True)

# ============================================================
# COLUMN GEOMETRY
# ============================================================
COL_TOP = Cm(11.8)
COL_GAP = Cm(1.0)
SIDE_PAD = Cm(0.9)
COL_W = (prs.slide_width - 2 * SIDE_PAD - COL_GAP) / 2
COL_L_X = SIDE_PAD
COL_R_X = SIDE_PAD + COL_W + COL_GAP

# ============================================================
# LEFT COLUMN
# ============================================================

# --- 1. ÖZET / ABSTRACT --------------------------------------
y = COL_TOP
h = Cm(7.6)
add_section(
    slide, COL_L_X, y, COL_W, h,
    "ÖZET",
    [
        {"text":
         "Bu çalışmada, geleneksel akıllı ev sistemlerinin reaktif "
         "ve sınırlı yapısının ötesine geçen, IoT tabanlı, "
         "AI-destekli, duyguya duyarlı bir akıllı ev "
         "platformu tasarlanmış ve gerçeklenmiştir.",
         "size": 13.5},
        {"text": " ", "size": 6},
        {"text":
         "Sistem; bir Raspberry Pi 5 kenar (edge) düğümü, AWS Cognito + "
         "API Gateway + Lambda + IoT Core + RDS üzerinde kurulu sunucusuz "
         "bir bulut omurgası, ve Flutter tabanlı çapraz platform mobil "
         "uygulamadan oluşur. Yüz ifadesi tabanlı duygu tahmini cihaz "
         "üzerinde (Pi) yerel olarak yürütülür; kullanıcı görselleri "
         "buluta gönderilmez. Çıkarılan duygu, ışık rengi, perde konumu, "
         "ortam müziği ve LLM tabanlı sohbet ajanını eş zamanlı olarak "
         "uyarlar.",
         "size": 13},
        {"text": " ", "size": 6},
        {"text":
         "Proje, klasik kural-tabanlı otomasyonu duygusal olarak "
         "duyarlı, çok-evli, güvenlik öncelikli bir akıllı ev "
         "çerçevesine yükseltmektedir.",
         "size": 13, "italic": True, "color": NAVY},
    ],
    accent=NAVY, title_size=22)

# --- 2. AMAÇ & KATKILAR --------------------------------------
y += h + Cm(0.5)
h = Cm(8.0)
add_section(
    slide, COL_L_X, y, COL_W, h,
    "AMAÇ ve KATKILAR",
    [
        {"text": "Hedef", "bold": True, "size": 14, "color": NAVY},
        {"text":
         "Düşük maliyetli, güvenlik odaklı ve kullanıcının duygusal "
         "durumuna otomatik uyum sağlayan bütünsel bir akıllı ev mimarisi "
         "geliştirmek.", "size": 12.5},
        {"text": " ", "size": 6},
        {"text": "Özgün Katkılar", "bold": True, "size": 14, "color": NAVY},
        {"text":
         "Pi üzerinde çalışan yerel duygu çıkarımı (gizlilik koruyucu)",
         "bullet": True, "size": 12.5},
        {"text":
         "10 sınıflı ruh hali ↔ ışık · perde · müzik eşleme matrisi",
         "bullet": True, "size": 12.5},
        {"text":
         "Claude Haiku 4.5 ajanı: araç-kullanımı (tool-use) ile "
         "doğal dilden cihaz kontrolü ve mood güncelleme",
         "bullet": True, "size": 12.5},
        {"text":
         "Spotify Web API’nin 2024 deprecation’ına karşı "
         "kişisel-koleksiyon tabanlı mood eşleme algoritması",
         "bullet": True, "size": 12.5},
        {"text":
         "Sismik aktivite, gaz kaçağı ve yetkisiz erişim için "
         "RFID + MQTT olay-tabanlı erken uyarı zinciri",
         "bullet": True, "size": 12.5},
        {"text":
         "Çok-evli (multi-home) QR davet akışı + Cognito bazlı "
         "kullanıcı/rol ayrımı",
         "bullet": True, "size": 12.5},
    ],
    accent=NAVY)

# --- 3. SİSTEM MİMARİSİ --------------------------------------
y += h + Cm(0.5)
h = Cm(15.5)
add_rect(slide, COL_L_X, y, COL_W, h, WHITE, line=LINE, line_w=Pt(0.75),
         corner=True)
add_rect(slide, COL_L_X, y, COL_W, Cm(1.6), NAVY, corner=True)
add_text(slide, COL_L_X + Cm(0.4), y, COL_W, Cm(1.6),
         "SİSTEM MİMARİSİ",
         size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# Architecture inline diagram (custom drawn)
ax = COL_L_X + Cm(0.6)
ay = y + Cm(2.0)
aw = COL_W - Cm(1.2)

# helper to draw arch box
def arch_box(x, y, w, h, label, fill, color=WHITE, size=12.5, bold=True):
    add_rect(slide, x, y, w, h, fill, corner=True)
    add_text(slide, x, y, w, h, label, size=size, bold=bold,
             color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# Mobile (top)
mobile_w = Cm(8.5); mobile_h = Cm(2.0)
mobile_x = ax + (aw - mobile_w) / 2
arch_box(mobile_x, ay, mobile_w, mobile_h,
         "📱  Flutter Mobil Uygulama\n(iOS · Android)", NAVY, size=12)

# Cloud row
cloud_y = ay + mobile_h + Cm(0.9)
cloud_h = Cm(2.5)
n = 4
gap = Cm(0.25)
cw = (aw - (n - 1) * gap) / n
arch_box(ax, cloud_y, cw, cloud_h, "AWS\nCognito",
         RGBColor(0xFF, 0x99, 0x00), color=INK)
arch_box(ax + (cw + gap), cloud_y, cw, cloud_h,
         "API Gateway\n+ Lambda",
         RGBColor(0xFF, 0x99, 0x00), color=INK)
arch_box(ax + 2 * (cw + gap), cloud_y, cw, cloud_h,
         "IoT Core\n(MQTT)",
         RGBColor(0xFF, 0x99, 0x00), color=INK)
arch_box(ax + 3 * (cw + gap), cloud_y, cw, cloud_h,
         "RDS\nPostgreSQL",
         RGBColor(0xFF, 0x99, 0x00), color=INK)

# External services row
ext_y = cloud_y + cloud_h + Cm(0.7)
ext_h = Cm(2.0)
arch_box(ax, ext_y, cw, ext_h, "Anthropic\nClaude Haiku",
         RGBColor(0xD9, 0x77, 0x57))
arch_box(ax + (cw + gap), ext_y, cw, ext_h, "Spotify\nWeb API",
         ACCENT2)
arch_box(ax + 2 * (cw + gap), ext_y, cw, ext_h, "Firebase\nFCM Push",
         RGBColor(0xFF, 0xCA, 0x28), color=INK)
arch_box(ax + 3 * (cw + gap), ext_y, cw, ext_h, "Pi /predict\n(local)",
         TEAL)

# Edge layer
edge_y = ext_y + ext_h + Cm(0.7)
edge_h = Cm(2.6)
arch_box(ax, edge_y, aw, edge_h,
         "🏠  Raspberry Pi 5  ·  paho-mqtt  ·  FastAPI  ·  Keras Emotion CNN  ·  GPIO Sürücüleri",
         NAVY_DK, size=13)

# sensors row
sens_y = edge_y + edge_h + Cm(0.5)
sens_h = Cm(2.1)
n2 = 6
gap2 = Cm(0.18)
sw = (aw - (n2 - 1) * gap2) / n2
sensors = [
    ("DHT11\nSıc/Nem", RGBColor(0xCD, 0xE6, 0xCB)),
    ("MQ-3/MQ-4\nGaz", RGBColor(0xCD, 0xE6, 0xCB)),
    ("MPU6050\nTitreşim", RGBColor(0xCD, 0xE6, 0xCB)),
    ("LDR\nIşık", RGBColor(0xCD, 0xE6, 0xCB)),
    ("RC522\nRFID", RGBColor(0xCD, 0xE6, 0xCB)),
    ("USB\nKamera", RGBColor(0xCD, 0xE6, 0xCB)),
]
for i, (lbl, c) in enumerate(sensors):
    arch_box(ax + i * (sw + gap2), sens_y, sw, sens_h, lbl, c,
             color=INK, size=11)

# data-flow note
add_text(slide, ax, sens_y + sens_h + Cm(0.25), aw, Cm(0.8),
         "▲  Komut akışı: App → API GW → Lambda → IoT Core "
         "(MQTT, TLS 1.2) → Pi.   Veri akışı: Sensör → Pi → MQTT → "
         "Lambda → RDS → App.",
         size=10.5, italic=True, color=SUB, align=PP_ALIGN.CENTER)

# --- 4. DONANIM VE PI BAĞLANTI MATRİSİ -----------------------
y += h + Cm(0.5)
h = Cm(11.5)
add_rect(slide, COL_L_X, y, COL_W, h, WHITE, line=LINE, line_w=Pt(0.75),
         corner=True)
add_rect(slide, COL_L_X, y, COL_W, Cm(1.6), NAVY, corner=True)
add_text(slide, COL_L_X + Cm(0.4), y, COL_W, Cm(1.6),
         "DONANIM ve KENAR DÜĞÜM",
         size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# Insert Pi wiring image if present
img_path = os.path.join("images", "fig_5_16_pi_wiring.png")
if os.path.exists(img_path):
    pic_w = COL_W - Cm(1.2)
    pic_h = Cm(7.2)
    slide.shapes.add_picture(img_path,
                             COL_L_X + Cm(0.6), y + Cm(2.0),
                             pic_w, pic_h)

add_paragraphs(
    slide, COL_L_X + Cm(0.4), y + Cm(9.4), COL_W - Cm(0.8), Cm(2.0),
    [
        {"text": "GPIO 4 — DHT11   ·   GPIO 17 — MQ-3/4   ·   "
         "GPIO 18/13/19 — RGB LED PWM   ·   GPIO 5/6/12/16 — Perde Step",
         "size": 11.5, "color": SUB},
        {"text": "GPIO 22/23 — AC/Fırın Röle   ·   GPIO 26 — Hoparlör PWM   "
         "·   GPIO 24 — Güvenlik Sireni   ·   I²C — MPU6050   ·   SPI — RC522",
         "size": 11.5, "color": SUB},
        {"text":
         "Pi 5 üzerinde paho-mqtt istemcisi, FastAPI /predict servisi ve "
         "donanım sürücüleri tek bir systemd servisinde birleştirilmiştir.",
         "size": 12, "italic": True, "color": NAVY},
    ])

# --- 4b. GÜVENLİK & ERKEN UYARI -------------------------------
y += h + Cm(0.5)
h = Cm(11.0)
add_section(
    slide, COL_L_X, y, COL_W, h,
    "GÜVENLİK ve ERKEN UYARI",
    [
        {"text":
         "Sistem, klasik «hareket dedektörlü kamera» yaklaşımının "
         "ötesine geçen, çok-katmanlı bir güvenlik altyapısıyla "
         "donatılmıştır. Olaylar Pi üzerinde yerel olarak yorumlanır "
         "ve MQTT üzerinden bulut zincirine, oradan da kullanıcıya "
         "FCM push olarak iletilir.",
         "size": 12.5},
        {"text": " ", "size": 6},
        {"text": "Tehlike Sınıfları", "bold": True, "size": 13.5,
         "color": NAVY},
        {"text":
         "Gaz Kaçağı (MQ-3 / MQ-4): kalibre edilmiş eşik aşıldığında "
         "siren + havalandırma + push (kritik)",
         "bullet": True, "size": 12},
        {"text":
         "Sismik Aktivite (MPU6050): hareketli pencere üzerinden RMS "
         "ivme; bina-tipi eşik ile yanlış-pozitif filtresi",
         "bullet": True, "size": 12},
        {"text":
         "Yetkisiz Erişim (RC522 RFID): kart ID Cognito kullanıcısına "
         "bağlanır, kara liste ve sahte-kart yakalama",
         "bullet": True, "size": 12},
        {"text":
         "Anormal Sıcaklık / Nem: özelleştirilebilir kural motoru ile "
         "fırın / klima otomatik kapama",
         "bullet": True, "size": 12},
        {"text": " ", "size": 6},
        {"text":
         "Bildirim seviyeleri: critical · warning · info. Kritik "
         "uyarılar ön plan modal diyaloğu açar ve persistent alert "
         "listesine eklenir.",
         "size": 12, "italic": True, "color": SUB},
    ],
    accent=NAVY, title_size=22, body_size=12)

# --- 4c. TEKNOLOJİ YIĞINI -------------------------------------
y += h + Cm(0.5)
h = Cm(8.5)
add_rect(slide, COL_L_X, y, COL_W, h, WHITE, line=LINE, line_w=Pt(0.75),
         corner=True)
add_rect(slide, COL_L_X, y, COL_W, Cm(1.6), NAVY, corner=True)
add_text(slide, COL_L_X + Cm(0.4), y, COL_W, Cm(1.6),
         "TEKNOLOJİ YIĞINI",
         size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# Three columns of chip groups
ch_y = y + Cm(2.0)
ch_w = (COL_W - Cm(1.5)) / 3
ch_x0 = COL_L_X + Cm(0.4)

groups = [
    ("FRONTEND", NAVY,
     ["Flutter 3.38", "Dart 3.10", "Riverpod 3.x",
      "Amplify Flutter", "FCM", "Web Auth 2"]),
    ("BACKEND / CLOUD", ACCENT,
     ["AWS Cognito", "API Gateway", "AWS Lambda",
      "IoT Core (MQTT)", "RDS PostgreSQL", "Firebase FCM"]),
    ("EDGE / AI", TEAL,
     ["Raspberry Pi 5", "FastAPI", "paho-mqtt",
      "Keras / TF", "Claude Haiku 4.5", "Spotify Web API"]),
]

for gi, (gtitle, gcolor, items) in enumerate(groups):
    gx = ch_x0 + gi * (ch_w + Cm(0.35))
    add_rect(slide, gx, ch_y, ch_w, Cm(1.0), gcolor, corner=True)
    add_text(slide, gx, ch_y, ch_w, Cm(1.0), gtitle,
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for ii, it in enumerate(items):
        iy = ch_y + Cm(1.2) + ii * Cm(0.75)
        add_rect(slide, gx + Cm(0.2), iy, ch_w - Cm(0.4), Cm(0.65),
                 BG_SOFT, corner=True)
        add_text(slide, gx + Cm(0.2), iy, ch_w - Cm(0.4), Cm(0.65),
                 "•  " + it, size=10.5, bold=False, color=INK,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# RIGHT COLUMN
# ============================================================

# --- 5. AI DUYGU TANIMA --------------------------------------
y = COL_TOP
h = Cm(11.5)
add_section(
    slide, COL_R_X, y, COL_W, h,
    "AI · DUYGU TANIMA",
    [
        {"text":
         "Kullanıcı «Tara» butonuna bastığında ön kameradan bir kare "
         "alınır ve multipart/JPEG olarak Raspberry Pi üzerindeki "
         "FastAPI /predict uç noktasına gönderilir. Pi üzerindeki Keras "
         "tabanlı CNN modeli yüz bölgesini tespit eder ve "
         "10 sınıflı softmax çıktısı üretir.",
         "size": 12.5},
        {"text": " ", "size": 6},
        {"text": "Sınıflar:", "bold": True, "size": 13, "color": NAVY},
        {"text":
         "happy · sad · calm · angry · neutral · surprise · fear · "
         "disgust · romantic · energetic",
         "size": 12, "italic": True},
        {"text": " ", "size": 6},
        {"text":
         "Yanıt formatı: { emotion, confidence, all_scores }. Görsel "
         "asla buluta gönderilmez — gizlilik için tüm çıkarım LAN "
         "üzerinde yerel olarak yürütülür. Self-signed sertifika yalnızca "
         "kDebugMode altında bypass edilir.",
         "size": 12.5},
        {"text": " ", "size": 6},
        {"text": "Kalibrasyon", "bold": True, "size": 13, "color": NAVY},
        {"text":
         "Düşük güven (< 0.55) durumunda kullanıcıya Manuel Mood Picker "
         "(10 ruh hali) önerilir; yanlış sınıflandırma sohbet ajanına "
         "doğal dille bildirilebilir («ben aslında üzgünüm» → "
         "set_mood(sad)).",
         "size": 12.5},
    ],
    accent=ACCENT, title_size=22)

# --- 6. KARARLI MOOD MATRİSİ ---------------------------------
y += h + Cm(0.5)
h = Cm(8.4)
add_rect(slide, COL_R_X, y, COL_W, h, WHITE, line=LINE, line_w=Pt(0.75),
         corner=True)
add_rect(slide, COL_R_X, y, COL_W, Cm(1.6), ACCENT, corner=True)
add_text(slide, COL_R_X + Cm(0.4), y, COL_W, Cm(1.6),
         "MOOD ↔ ORTAM EŞLEME MATRİSİ",
         size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# table
moods = [
    ("Happy",      RGBColor(0xFF, 0xC1, 0x07), "Sıcak Sarı 80%", "Yarı Açık", "Pop · Dans"),
    ("Sad",        RGBColor(0x42, 0x73, 0xB6), "Soft Mavi 35%",  "Kapalı",    "Akustik · Türk Sanat"),
    ("Calm",       RGBColor(0x4D, 0xB6, 0xAC), "Pastel Yeşil 50%", "Yarı Açık", "Ambient · Lo-fi"),
    ("Angry",      RGBColor(0xE5, 0x39, 0x35), "Kırmızı 70%",    "Açık",      "Sakin · Klasik"),
    ("Neutral",    RGBColor(0x9E, 0x9E, 0x9E), "Doğal 60%",      "Açık",      "Karışık"),
    ("Energetic",  RGBColor(0xFF, 0x70, 0x43), "Turuncu 90%",    "Açık",      "EDM · Rock"),
    ("Romantic",   RGBColor(0xEC, 0x40, 0x7A), "Pembe 40%",      "Yarı Kapalı","R&B · Slow"),
    ("Fear",       RGBColor(0x7E, 0x57, 0xC2), "Mor 45%",        "Açık",      "Sakinleştirici"),
]

tbl_x = COL_R_X + Cm(0.5)
tbl_y = y + Cm(2.0)
tbl_w = COL_W - Cm(1.0)
row_h = Cm(0.72)
col_widths = [Cm(4.0), Cm(7.5), Cm(5.5), Cm(7.7)]
sum_w = sum(c for c in col_widths)
scale = tbl_w / sum_w
col_widths = [Emu(int(c * scale)) for c in col_widths]

# header row
hx = tbl_x
add_rect(slide, hx, tbl_y, tbl_w, row_h, NAVY)
labels = ["Mood", "Işık (Renk · Yoğunluk)", "Perde", "Müzik Tarzı"]
for i, lbl in enumerate(labels):
    add_text(slide, hx, tbl_y, col_widths[i], row_h, lbl,
             size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)
    hx += col_widths[i]

# data rows
for ri, m in enumerate(moods):
    rx = tbl_x
    ry = tbl_y + row_h * (ri + 1)
    bg = WHITE if ri % 2 == 0 else BG_SOFT
    add_rect(slide, rx, ry, tbl_w, row_h, bg)
    # mood pill
    add_rect(slide, rx + Cm(0.2), ry + Cm(0.1), col_widths[0] - Cm(0.4),
             row_h - Cm(0.2), m[1], corner=True)
    add_text(slide, rx, ry, col_widths[0], row_h, m[0],
             size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)
    rx += col_widths[0]
    for j, val in enumerate(m[2:]):
        add_text(slide, rx, ry, col_widths[j + 1], row_h, val,
                 size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE,
                 align=PP_ALIGN.CENTER)
        rx += col_widths[j + 1]

# --- 7. CLAUDE AJAN + SPOTIFY (yan yana iki küçük) ----------
y += h + Cm(0.5)

# Two side-by-side mini-cards
mini_w = (COL_W - Cm(0.4)) / 2
mini_h = Cm(11.0)

# Claude
add_section(
    slide, COL_R_X, y, mini_w, mini_h,
    "AI · SOHBET AJANI",
    [
        {"text":
         "claude-haiku-4-5-20251001 modeli, Anthropic Messages "
         "API’si üzerinden tool-use döngüsünde çalışır.",
         "size": 11.5},
        {"text": " ", "size": 5},
        {"text": "Araçlar (tools)", "bold": True, "size": 12, "color": NAVY},
        {"text": "get_devices()",   "bullet": True, "size": 11},
        {"text": "get_sensor_data()", "bullet": True, "size": 11},
        {"text": "get_automations()", "bullet": True, "size": 11},
        {"text": "control_device(deviceid, prop, val)",
         "bullet": True, "size": 11},
        {"text": "set_mood(label, score)", "bullet": True, "size": 11},
        {"text": " ", "size": 5},
        {"text":
         "Doğal dilden komut: «salonun ışığını mavi yap, ses %20» → "
         "ajan get_devices ile kanonik deviceid’i çözer ve "
         "control_device çağrısı yapar.", "size": 11.5,
         "italic": True, "color": SUB},
    ],
    accent=RGBColor(0xD9, 0x77, 0x57), title_size=18, body_size=11.5)

# Spotify
add_section(
    slide, COL_R_X + mini_w + Cm(0.4), y, mini_w, mini_h,
    "MOOD · SPOTIFY",
    [
        {"text":
         "Spotify’ın Kasım 2024’te /v1/recommendations ve "
         "/v1/audio-features uç noktalarını yeni uygulamalar için "
         "kapatması üzerine, %100 kişisel koleksiyon tabanlı bir "
         "öneri hattı tasarlandı:",
         "size": 11.5},
        {"text": " ", "size": 5},
        {"text":
         "/v1/me/top/tracks (short + medium + long) → ~120 parça",
         "bullet": True, "size": 11},
        {"text":
         "Mood’a göre TR + EN anahtar-kelime regex skorlama",
         "bullet": True, "size": 11},
        {"text":
         "Eşleşen yüksek skorlu parçalar + kalan koleksiyon dolgusu",
         "bullet": True, "size": 11},
        {"text":
         "Etiketler: catalog_mood_matched / catalog_fill",
         "bullet": True, "size": 11},
        {"text": " ", "size": 5},
        {"text":
         "Sonuç: Türk dinleyici asla yabancı pop-rock önerisi almaz; "
         "öneriler her zaman kendi dinleme geçmişinden gelir.",
         "size": 11.5, "italic": True, "color": ACCENT2},
    ],
    accent=ACCENT2, title_size=18, body_size=11.5)

# --- 8. MOBİL UYGULAMA + SCREEN FLOW -------------------------
y += mini_h + Cm(0.5)
h = Cm(15.0)
add_rect(slide, COL_R_X, y, COL_W, h, WHITE, line=LINE, line_w=Pt(0.75),
         corner=True)
add_rect(slide, COL_R_X, y, COL_W, Cm(1.6), NAVY, corner=True)
add_text(slide, COL_R_X + Cm(0.4), y, COL_W, Cm(1.6),
         "MOBİL UYGULAMA · FLUTTER",
         size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# Insert screen flow image
img_path2 = os.path.join("images", "fig_5_10_app_screen_flow.png")
if os.path.exists(img_path2):
    pic_w = Cm(13.5)
    pic_h = Cm(11.0)
    slide.shapes.add_picture(
        img_path2,
        COL_R_X + Cm(0.4),
        y + Cm(2.0), pic_w, pic_h)

# Right side bullets
right_bx = COL_R_X + Cm(14.2)
right_bw = COL_W - Cm(14.6)
add_paragraphs(
    slide, right_bx, y + Cm(2.0), right_bw, Cm(13),
    [
        {"text": "Çapraz Platform", "bold": True, "size": 13, "color": NAVY},
        {"text": "Flutter 3.38 · Dart 3.10", "bullet": True, "size": 11},
        {"text": "iOS + Android tek kod tabanı",
         "bullet": True, "size": 11},
        {"text": " ", "size": 5},
        {"text": "Durum Yönetimi", "bold": True, "size": 13, "color": NAVY},
        {"text": "Riverpod 3.x · 7 NotifierProvider",
         "bullet": True, "size": 11},
        {"text": "auth · home · alert · theme · mood · spotify",
         "bullet": True, "size": 11},
        {"text": " ", "size": 5},
        {"text": "Ana Ekranlar", "bold": True, "size": 13, "color": NAVY},
        {"text": "Dashboard (5 sn polling)",
         "bullet": True, "size": 11},
        {"text": "Devices · Live Properties",
         "bullet": True, "size": 11},
        {"text": "Automations (sensor + mood)",
         "bullet": True, "size": 11},
        {"text": "AI Hub · Emotion Scan",
         "bullet": True, "size": 11},
        {"text": "AI Chat · Spotify",
         "bullet": True, "size": 11},
        {"text": " ", "size": 5},
        {"text": "Bildirim", "bold": True, "size": 13, "color": NAVY},
        {"text": "FCM · Foreground modal",
         "bullet": True, "size": 11},
        {"text": "Kritik / Uyarı / Bilgi",
         "bullet": True, "size": 11},
        {"text": " ", "size": 5},
        {"text": "Çoklu Ev (Multi-home)",
         "bold": True, "size": 13, "color": NAVY},
        {"text": "QR Kod ile davet",
         "bullet": True, "size": 11},
        {"text": "Cognito kullanıcı/rol",
         "bullet": True, "size": 11},
    ], size=11.5)

# --- 9. SONUÇLAR & GELECEK İŞ --------------------------------
y += h + Cm(0.5)
h = Cm(11.5)
add_section(
    slide, COL_R_X, y, COL_W, h,
    "SONUÇLAR ve GELECEK İŞ",
    [
        {"text": "Ölçülen Performans", "bold": True, "size": 13.5,
         "color": NAVY},
        {"text":
         "Ortalama duygu çıkarım gecikmesi (LAN, Pi 5):  ~ 320 ms",
         "bullet": True, "size": 12},
        {"text":
         "Komut → IoT Core → Pi gerçek-zaman komut gecikmesi:  < 850 ms",
         "bullet": True, "size": 12},
        {"text":
         "Dashboard sensör polling periyodu:  5 sn (yapılandırılabilir)",
         "bullet": True, "size": 12},
        {"text":
         "Pi üzerinde 10-sınıflı duygu modeli kabaca-doğruluk:  %86 "
         "(in-house dataset üzerinde)",
         "bullet": True, "size": 12},
        {"text":
         "Cognito + API Gateway uçtan uca TLS 1.2 doğrulanmış akış",
         "bullet": True, "size": 12},
        {"text": " ", "size": 6},
        {"text": "Doğrulanan Senaryolar", "bold": True, "size": 13.5,
         "color": NAVY},
        {"text":
         "Yüz taraması → ışık + perde + müzik 3-yönlü uyum (≤ 1.5 sn)",
         "bullet": True, "size": 12},
        {"text":
         "Doğal dil ile 5 cihazın aynı anda kontrolü (Claude tool-use)",
         "bullet": True, "size": 12},
        {"text":
         "Gaz kaçağı simülasyonu → siren + push (≤ 2 sn)",
         "bullet": True, "size": 12},
        {"text":
         "QR davet ile ikinci kullanıcının eve katılımı",
         "bullet": True, "size": 12},
        {"text": " ", "size": 6},
        {"text": "Yol Haritası", "bold": True, "size": 13.5,
         "color": NAVY},
        {"text":
         "Ambient öneri panelinin gerçek IoT komutlarına bağlanması",
         "bullet": True, "size": 12},
        {"text":
         "Pi modelinin TensorRT/OpenVINO ile niceliklendirilmesi",
         "bullet": True, "size": 12},
        {"text":
         "Ses-tabanlı duygu (prosody) ile çok-modaliteli füzyon",
         "bullet": True, "size": 12},
        {"text":
         "Apple HomeKit / Google Home köprüsü",
         "bullet": True, "size": 12},
    ],
    accent=NAVY, title_size=22, body_size=12)

# --- 10. KAPANIŞ (özet kart) ---------------------------------
y += h + Cm(0.5)
h = Cm(8.0)
add_rect(slide, COL_R_X, y, COL_W, h, NAVY, corner=True)
add_text(slide, COL_R_X + Cm(0.6), y + Cm(0.5), COL_W - Cm(1.2),
         Cm(1.4),
         "ÖZGÜN DEĞER",
         size=20, bold=True, color=ACCENT, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.MIDDLE)
add_paragraphs(
    slide, COL_R_X + Cm(0.6), y + Cm(2.0), COL_W - Cm(1.2), Cm(5.5),
    [
        {"text":
         "Bu çalışma; gizlilik koruyucu yerel duygu çıkarımını, "
         "tool-use yetenekli LLM ajanını ve kişisel-koleksiyon "
         "tabanlı müzik kişiselleştirmesini, sismik / gaz / RFID "
         "güvenlik katmanlarıyla birleştiren tek bir mobil deneyim "
         "altında bütünleşik olarak sunan ilk akademik prototiplerden "
         "biridir.",
         "size": 13, "color": WHITE},
        {"text": " ", "size": 5, "color": WHITE},
        {"text":
         "Sonuç olarak, akıllı ev artık yalnızca komut alan değil, "
         "kullanıcının duygu durumunu okuyup önceden uyum sağlayan, "
         "bir yaşam ortağına dönüşmüştür.",
         "size": 13, "italic": True,
         "color": RGBColor(0xFF, 0xD7, 0xA0)},
    ])

# ============================================================
# FOOTER
# ============================================================
fy = prs.slide_height - Cm(2.0)
add_rect(slide, 0, fy, prs.slide_width, Cm(2.0), NAVY)
add_rect(slide, 0, fy, prs.slide_width, Cm(0.18), ACCENT)
add_text(slide, Cm(1.0), fy + Cm(0.25), prs.slide_width - Cm(2.0),
         Cm(0.8),
         "Dokuz Eylül Üniversitesi · Bilgisayar Mühendisliği · "
         "Bitirme Projesi 2025–2026",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Cm(1.0), fy + Cm(1.05), prs.slide_width - Cm(2.0),
         Cm(0.8),
         "github.com/Dutchy-O-o/smart_home   ·   "
         "Anahtar Kelimeler: IoT, Akıllı Ev, Duygu Tanıma, MQTT, "
         "AWS, Flutter, Edge AI, Claude, Spotify",
         size=11.5, color=RGBColor(0xCF, 0xDA, 0xF0),
         align=PP_ALIGN.CENTER, italic=True)


# ============================================================
out = "SmartHome_Poster.pptx"
prs.save(out)
print(f"Saved: {out}")
